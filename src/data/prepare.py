"""Document ingestion and chunking utilities."""

from __future__ import annotations

import json
import re
import subprocess
import sys
from collections import deque
from pathlib import Path
from typing import Deque, Iterable, List, Tuple

import html2text
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from config import BaseModelPaths, FolderPaths


def caption_image_file(image_path: Path) -> str:
    """Generate a caption for an image using mlx_vlm if available."""

    base_paths = BaseModelPaths()
    model_choice_str = str(base_paths.get_model_path("mlx_image"))

    prompt = "USER: <image>\nDescribe this image: \nASSISTANT:"
    print(f"Captioning image {image_path}")
    command = [
        sys.executable,
        "-m",
        "mlx_vlm.generate",
        "--model",
        model_choice_str,
        "--prompt",
        prompt,
        "--image",
        str(image_path),
        "--max-tokens",
        "256",
        "--temp",
        "0.1",
    ]

    try:
        result = subprocess.run(command, capture_output=True, text=True, check=False)
    except FileNotFoundError:
        print("mlx_vlm.generate not found; skipping caption generation.")
        return ""

    if result.returncode != 0:
        print("Error in generating caption:", result.stderr)
        return ""

    try:
        output = json.loads(result.stdout)
    except json.JSONDecodeError:
        print("Invalid caption output:", result.stdout)
        return ""

    caption = output.get("caption", "")
    print("Caption:", caption)
    return caption


def clean_text(line: str) -> str | None:
    line = line.strip()
    if len(re.findall(r"\w+", line)) < 3 or line.count("\t") > 1:
        return None
    return line


def read_txt_file(file_path: Path, _: int, __) -> str:
    return file_path.read_text()


def handle_image(
    image_bytes: bytes,
    folders: FolderPaths,
    file_index: int,
    page_number: int,
    image_index: int,
    suffix: str,
    args,
) -> str:
    image_filename = f"output_file_{file_index}{suffix}_page_{page_number}_image_{image_index}.jpg"
    image_path = folders.image_folder / image_filename
    image_path.write_bytes(image_bytes)
    caption = caption_image_file(image_path) if args.images else ""
    print(f"Image saved to {image_path}")
    return f"\nImage Caption: {caption}\n" if caption else ""


def read_pdf_file(file_path: Path, file_index: int, args) -> str:
    print(f"Reading PDF file with images enabled: {args.images}")
    text = ""
    reader = PdfReader(str(file_path))
    image_index = 0
    folders = FolderPaths(args)
    for page_number, page in enumerate(reader.pages, start=1):
        extracted = page.extract_text()
        if extracted:
            text += extracted
        if args.images:
            images = getattr(page, "images", [])
            print(f"Found {len(images)} image(s) on page {page_number}")
            for image in images:
                image_bytes = getattr(image, "data", None)
                if not image_bytes:
                    continue
                text += handle_image(
                    image_bytes,
                    folders,
                    file_index,
                    page_number,
                    image_index,
                    "_pdf",
                    args,
                )
                image_index += 1
    return text


def read_pptx_file(file_path: Path, file_index: int, args) -> str:
    folders = FolderPaths(args)
    text = ""
    prs = Presentation(str(file_path))
    image_index = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        if slide.shapes.title:
            text += slide.shapes.title.text + "\n\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    bullet_char = "• " if paragraph.level > 0 else ""
                    text += bullet_char + paragraph.text + "\n"
                text += "\n"
            elif hasattr(shape, "image") and args.images:
                image_bytes = shape.image.blob
                image_filename = (
                    f"output_file_{file_index}_pptx_slide_{slide_number}_image_{image_index}.jpg"
                )
                image_path = folders.image_folder / image_filename
                image_path.write_bytes(image_bytes)
                caption = caption_image_file(image_path) if args.images else ""
                if caption:
                    text += f"\nImage Caption: {caption}\n"
                print(f"Image saved to {image_path}")
                image_index += 1
        text += "-----\n"
    return text


def read_html_file(file_path: Path, _: int, __) -> str:
    html = file_path.read_text()
    parser = html2text.HTML2Text()
    parser.ignore_links = True
    return parser.handle(html)


def read_docx_file(file_path: Path, _: int, __) -> str:
    document = Document(str(file_path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


FILE_READERS = {
    "pdf": read_pdf_file,
    "pptx": read_pptx_file,
    "txt": read_txt_file,
    "html": read_html_file,
    "docx": read_docx_file,
}


def process_file(file_path: Path, file_type: str, file_index: int, args) -> str | None:
    reader_func = FILE_READERS.get(file_type)
    if reader_func:
        return reader_func(file_path, file_index, args)
    return None


def _iter_clean_lines(text: str) -> Iterable[str]:
    for line in text.split("\n"):
        cleaned = clean_text(line)
        if cleaned:
            yield cleaned


def _write_chunk(output_folder: Path, chunk_index: int, word_limit: int, words: List[str]) -> Path:
    output_file = f"output_chunk_{chunk_index}_{word_limit}.txt"
    output_path = output_folder / output_file
    output_path.write_text(" ".join(words) + "\n")
    return output_path


def concatenate_files(
    folders: FolderPaths,
    file_queue: Deque[Tuple[Path, str, int]],
    args,
) -> List[Path]:
    output_folder = folders.prepared_data_folder
    chunk_index = 1
    current_words: List[str] = []
    chunk_paths: List[Path] = []
    word_limit = max(args.word_limit, 1)
    overlap_words = int(word_limit * args.overlap)

    for file_path, file_type, file_index in file_queue:
        file_text = process_file(file_path, file_type, file_index, args)
        if not file_text:
            continue

        for line in _iter_clean_lines(file_text):
            words = line.split()
            current_words.extend(words)

            while len(current_words) >= word_limit:
                chunk_words = current_words[:word_limit]
                chunk_paths.append(_write_chunk(output_folder, chunk_index, word_limit, chunk_words))
                chunk_index += 1
                current_words = current_words[word_limit - overlap_words if overlap_words else word_limit:]

    if current_words:
        chunk_paths.append(_write_chunk(output_folder, chunk_index, word_limit, current_words))

    return chunk_paths


def file_processor(folders: FolderPaths, args) -> List[Path]:
    documents_folder = folders.documents_folder
    print(f"Appending data from {documents_folder}")

    supported_types = {"pdf", "pptx", "txt", "html", "docx"}
    try:
        files = [
            (path, path.suffix.lstrip(".").lower())
            for path in documents_folder.iterdir()
            if path.is_file() and path.suffix.lstrip(".").lower() in supported_types
        ]
    except FileNotFoundError:
        print(f"Documents folder {documents_folder} not found.")
        return []

    if not files:
        print("No files found in the input folder.")
        return []

    sorted_files = sorted(files, key=lambda item: item[0].name)
    file_queue: Deque[Tuple[Path, str, int]] = deque(
        (path, ext, index)
        for index, (path, ext) in enumerate(sorted_files, start=1)
    )
    print(f"Processing {len(file_queue)} files.")
    return concatenate_files(folders, file_queue, args)
