# prepare_data.py (refactored)

import os
from pptx import Presentation
from docx import Document
import html2text
import re
from setup_and_parse import FolderPaths, BaseModelPaths
from PIL import Image
from collections import deque
from stqdm import stqdm
from PyPDF2 import PdfReader

def caption_image_file(image_path):
    import subprocess
    import json
    import sys  # Import sys to use sys.executable
    
    # If image is lower than 256x256 width or height, skip image
    # image = Image.open(image_path)
    # if image.width < 256 or image.height < 256:
    #     return "" # Return nothing if image is too small
    
    base_paths = BaseModelPaths()
    model_choice_str = str(base_paths.get_model_path('mlx_image'))
    
    prompt = str("USER: <image>\nDescribe this image: \nASSISTANT:")
    max_tokens = 256
    temperature = 0.1
    
    print(f"Captioning image {image_path}")
    # Construct the command to run the model as a subprocess using sys.executable
    command = [
        sys.executable, "-m", "mlx_vlm.generate",
        "--model", model_choice_str,
        "--prompt", prompt,
        "--image", image_path,
        "--max-tokens", str(max_tokens),
        "--temp", str(temperature)
    ]
    
    # Execute the command and capture the output
    result = subprocess.run(command, capture_output=True, text=True)
    if result.returncode != 0:
        print("Error in generating caption:", result.stderr)
        return ""
    
    # Assuming the output is JSON formatted
    output = json.loads(result.stdout)
    caption = output.get("caption", "")
    print("Caption: ", caption)
    return caption

def clean_text(line):
    line = line.strip()
    if len(re.findall(r'\w+', line)) < 3 or line.count("\t") > 1:
        return None
    return line

def read_txt_file(file_path, file_index, args):
    with open(file_path, 'r') as f:
        return f.read()

def handle_image(doc, img, file_index, page_number, image_index, file_type, args):
    folders = FolderPaths(args)
    base_image = doc.extract_image(img[0])
    image_bytes = base_image["image"]
    image_filename = f"output_file_{file_index}{file_type}_page_{page_number}_image_{image_index}.jpg"
    image_path = os.path.join(folders.image_folder, image_filename)
    with open(image_path, 'wb') as img_file:
        img_file.write(image_bytes)
    caption = caption_image_file(image_path)
    print(f"Image saved to {image_path}")
    return f"\nImage Caption: {caption}\n"

def read_pdf_file(file_path, file_index, args):
    print(f"Reading PDF file with images enabled: {args.images}")
    text = ""
    reader = PdfReader(file_path)
    image_index = 0
    for page_number, page in enumerate(reader.pages, start=1):
        text += page.extract_text() if page.extract_text() else ""
        if args.images:
            print(f"Found {len(page.images)} image(s) on page {page_number}")
            for img_key in page.images:
                if isinstance(img_key, str):
                    # Top-level image
                    image = page.images[img_key]
                elif isinstance(img_key, tuple):
                    # Image within an XObject Form
                    image = page.images[img_key[0]][img_key[1]]
                else:
                    # Inline image
                    image = page.images[img_key]
                
                # Process each image using the existing handle_image function
                text += handle_image(reader, image, file_index, page_number, image_index, "_pdf", args)
                image_index += 1
    return text

def read_pptx_file(file_path, file_index, args):
    folders = FolderPaths(args)
    text = ""
    prs = Presentation(file_path)
    image_index = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        if slide.shapes.title:
            text += slide.shapes.title.text + "\n\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    bullet_char = "• " if paragraph.level > 0 else ""
                    text += bullet_char + paragraph.text + "\n"
                text += "\n"
            elif hasattr(shape, "image") and args.images:
                image_bytes = shape.image.blob
                image_filename = f"output_file_{file_index}_pptx_slide_{slide_number}_image_{image_index}.jpg"
                image_path = os.path.join(folders.image_folder, image_filename)
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_bytes)
                caption = caption_image_file(image_path)
                text += f"\nImage Caption: {caption}\n"
                print(f"Image saved to {image_path}")
                image_index += 1
        text += "-----\n"
    return text

def read_html_file(file_path, file_index, args):
    with open(file_path, 'r') as f:
        html = f.read()
    h = html2text.HTML2Text()
    h.ignore_links = True
    text = h.handle(html)

    if args.images:
        pass
        # Fix image captioning 
    return text

def read_docx_file(file_path, file_index, args):
    text = ""
    document = Document(file_path)
    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"
    
    if args.images:
        # Fix image captioning 
        pass
    return text


FILE_READERS = {
    'pdf': read_pdf_file,
    'pptx': read_pptx_file,
    'txt': read_txt_file,
    'html': read_html_file,
    'docx': read_docx_file
}

def process_file(file_path, file_type, file_index, args):
    reader_func = FILE_READERS.get(file_type)
    if reader_func:
        return reader_func(file_path, file_index, args)
    else:
        return None

def concatenate_files(folders, file_queue, args, all_data, file_contents):
    output_folder = folders.prepared_data_folder
    file_index = 1
    chunk_index = 1
    current_chunk = []
    current_chunk_word_count = 0
    overlap_word_count = int(args.word_limit * (args.overlap / 100))  # Calculate the overlap word count
    document_count = 0

    for current_file, file_type in file_queue:
        file_path = os.path.join(folders.documents_folder, current_file)
        file_text = process_file(file_path, file_type, file_index, args)
        if file_text:
            document_count += 1
            lines = file_text.split('\n')
            for line in lines:
                cleaned_line = clean_text(line)
                if cleaned_line:
                    words = cleaned_line.split()
                    current_chunk.extend(words)
                    current_chunk_word_count += len(words)

                    while current_chunk_word_count >= args.word_limit:
                        # Write current chunk to file with overlap
                        output_file = f"output_file_{file_index}_{chunk_index}_{args.word_limit}.txt"
                        chunk_words = current_chunk[:args.word_limit]
                        with open(os.path.join(output_folder, output_file), 'w') as out_file:
                            out_file.write(' '.join(chunk_words) + '\n')
                        all_data.extend(chunk_words)
                        file_contents.extend(chunk_words)

                        # Reset for new chunk with overlap
                        current_chunk = current_chunk[args.word_limit - overlap_word_count:]
                        current_chunk_word_count = len(current_chunk)
                        chunk_index += 1

        file_index += 1  # Increment file_index per file to avoid overwriting in different file types

    # Write the remaining chunk after processing all files
    if current_chunk:
        output_file = f"output_file_{file_index}_{chunk_index}_{args.word_limit}.txt"
        with open(os.path.join(output_folder, output_file), 'w') as out_file:
            out_file.write(' '.join(current_chunk) + '\n')
        all_data.extend(current_chunk)
        file_contents.extend(current_chunk)
        chunk_index += 1
        
    return chunk_index - 1

def file_processor(folders, args):
    documents_folder = folders.documents_folder
    print(f"Appending data from {documents_folder}")
    try:
        file_types = ['pdf', 'pptx', 'txt', 'html', 'docx']
        files = [(os.path.join(documents_folder, file), file.split('.')[-1]) for file in os.listdir(documents_folder) if file.split('.')[-1] in file_types]
        file_queue = deque(sorted(files))
        print(f"Processing {len(file_queue)} files.")
    except Exception as e:
        print(f"Error reading input folder {documents_folder}: {e}")
        return [], []

    if not file_queue:
        print("No files found in the input folder.")
        return [], []

    all_data = []
    file_contents = []
    concatenate_files(folders, file_queue, args, all_data, file_contents)
    return all_data, file_contents